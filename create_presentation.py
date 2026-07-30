import sys
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("python-pptx not found. Installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Clean Premium Farming Theme
    PRIMARY_COLOR = RGBColor(76, 175, 80)     # Vibrant Green (#4CAF50)
    DARK_COLOR = RGBColor(30, 41, 59)         # Slate Dark (#1E293B)
    LIGHT_BG = RGBColor(248, 250, 252)        # Very light grey (#F8FAFC)
    SECONDARY_COLOR = RGBColor(14, 165, 233)  # Sky Blue (#0EA5E9)
    TEXT_COLOR = RGBColor(51, 65, 85)         # Slate Gray (#334155)

    # Helper function to style text frame
    def add_slide_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = DARK_COLOR
        return title_box

    # Slide 1: Title Slide (Custom Minimal Elegant Layout)
    slide_layout = prs.slide_layouts[6] # Blank Layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Large colored title block or textbox
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.833), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "🌱 KrishiDrishti"
    p.font.name = 'Arial'
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Intelligent Farming Assistant"
    p2.font.name = 'Arial'
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = DARK_COLOR
    p2.space_before = Pt(10)
    p2.alignment = PP_ALIGN.LEFT

    p3 = tf.add_paragraph()
    p3.text = "A Production-Ready Full-Stack Web Application for Modern Agriculture"
    p3.font.name = 'Arial'
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_COLOR
    p3.space_before = Pt(30)
    p3.alignment = PP_ALIGN.LEFT

    # Slide 2: Project Overview & Objectives
    slide = prs.slides.add_slide(slide_layout)
    add_slide_title(slide, "Project Overview")
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("Core Vision", "KrishiDrishti is a modern full-stack web application designed to bridge the digital divide for small and marginal farmers, providing real-time data, localized advisories, and policy assistance."),
        ("Key Goals", "Empower agricultural communities through reliable real-time weather analytics, mandi market price monitoring, and simplified access to government support programs."),
        ("User Accessibility", "Bilingual interface supporting both English and Hindi, built with a clean, high-performance visual dashboard that works efficiently on all devices."),
        ("Administrative Control", "Centralized admin dashboard equipped with real-time broadcast capabilities (via WebSockets) to push critical warnings (e.g., pests, extreme weather alerts) to active farmers instantly.")
    ]
    
    for i, (title, desc) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run_title = p.add_run()
        run_title.text = f"•  {title}: "
        run_title.font.bold = True
        run_title.font.size = Pt(18)
        run_title.font.color.rgb = PRIMARY_COLOR
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(16)
        run_desc.font.color.rgb = TEXT_COLOR

    # Slide 3: Technical Architecture
    slide = prs.slides.add_slide(slide_layout)
    add_slide_title(slide, "Technical Architecture")
    
    # Multi-column description
    col1 = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Frontend Architecture (Vite + React)"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_after = Pt(10)
    
    feats_fe = [
        "Vite & React SPA: For rapid UI rendering and state management.",
        "Tailwind CSS: Modern premium styling system with customized green-focused layout.",
        "Lucide Icons & Recharts: Elegant visualization of market trends and farming statistics.",
        "React Context API: Manages user session state and real-time socket connections globally."
    ]
    for feat in feats_fe:
        p2 = tf1.add_paragraph()
        p2.text = f"✔  {feat}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_COLOR
        p2.space_after = Pt(6)

    col2 = slide.shapes.add_textbox(Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = col2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "Backend & Database (Node.js + MongoDB)"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_after = Pt(10)
    
    feats_be = [
        "Express.js REST API: Clean MVC architecture separation of controllers, schemas, and routes.",
        "MongoDB Atlas (Mongoose): High-performance document database storing structured schemas.",
        "Socket.io: Enables instant, real-time weather & advisory notification push.",
        "JWT Authentication: Secure user session authorization with token verification middlewares."
    ]
    for feat in feats_be:
        p2 = tf2.add_paragraph()
        p2.text = f"✔  {feat}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_COLOR
        p2.space_after = Pt(6)

    # Slide 4: Farmer Workspace & Core Features
    slide = prs.slides.add_slide(slide_layout)
    add_slide_title(slide, "Farmer Workspace & Core Features")
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    features = [
        ("Dynamic Dashboard", "Provides a welcoming layout displaying customized localized content, unread alert status, and key agriculture metrics."),
        ("Localized Weather & Farming Advisories", "Dynamically determines temperature, wind, and humidity indicators and generates immediate, action-oriented agricultural tips (e.g., 'ideal time for spraying', 'postpone watering')."),
        ("Smart Government Schemes Directory", "Integrates complex scheme information categorized by utility (Financial, Irrigation, Insurance, Input Subsidies) and extracts clear benefit and eligibility rules for the farmer."),
        ("Bilingual Settings Panel", "Allows full translation toggling between English and Hindi, updating profile configurations in real-time across database sessions.")
    ]
    
    for i, (title, desc) in enumerate(features):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run_title = p.add_run()
        run_title.text = f"•  {title}: "
        run_title.font.bold = True
        run_title.font.size = Pt(18)
        run_title.font.color.rgb = SECONDARY_COLOR
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(16)
        run_desc.font.color.rgb = TEXT_COLOR

    # Slide 5: Admin Control Panel
    slide = prs.slides.add_slide(slide_layout)
    add_slide_title(slide, "Admin Control Panel & WebSockets")
    
    col1 = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "Centralized Admin Controls"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_after = Pt(10)
    
    admin_points = [
        "Farmers Registry: Detailed overview table of all registered farmers, searchable by location/state.",
        "System Statistics: Live monitoring of active system metrics, registrations, and database collections.",
        "Secure Role Guards: Route-level middleware ensures only authorized administrative profiles access administrative APIs."
    ]
    for point in admin_points:
        p2 = tf1.add_paragraph()
        p2.text = f"➤  {point}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_COLOR
        p2.space_after = Pt(10)

    col2 = slide.shapes.add_textbox(Inches(6.98), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = col2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "Socket.io Broadcast Warnings"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_COLOR
    p.space_after = Pt(10)
    
    socket_points = [
        "Real-Time Push: Broadcast warnings are dispatched instantly to all active connected user sessions without reloading.",
        "Category Filtering: Supports customized warnings for specific issues like Weather warnings, general notices, or price changes.",
        "Visual Highlights: Interactive badge system highlights and updates unread counts dynamically on the client navbar."
    ]
    for point in socket_points:
        p2 = tf2.add_paragraph()
        p2.text = f"➤  {point}"
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_COLOR
        p2.space_after = Pt(10)

    # Slide 6: Database Schema & Design
    slide = prs.slides.add_slide(slide_layout)
    add_slide_title(slide, "Mongoose Schemas & Data Model")
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    schemas = [
        ("User (Farmer)", "Fields: name, email, password (hashed), phone, location, state, preferredLanguage ('en'|'hi'), resetOTP, resetOTPExpires."),
        ("Admin Schema", "Fields: name, email, password, role ('admin'). Seeded automatically on application start if none exists."),
        ("WeatherHistory", "Fields: location, temperature, humidity, windSpeed, rainForecast, recommendations. Logs queries to run analytics."),
        ("Notification Schema", "Fields: user (null for global broadcast), type ('weather'|'scheme'|'price'|'general'), title, message, isRead."),
        ("Government Scheme", "Fields: title, description, eligibility, benefits, applyLink, category ('Financial Support'|'Crop Insurance'|'Irrigation & Power'|'Subsidies & Inputs').")
    ]
    
    for i, (name, details) in enumerate(schemas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run_name = p.add_run()
        run_name.text = f"•  {name}: "
        run_name.font.bold = True
        run_name.font.size = Pt(15)
        run_name.font.color.rgb = PRIMARY_COLOR
        
        run_details = p.add_run()
        run_details.text = details
        run_details.font.size = Pt(14)
        run_details.font.color.rgb = TEXT_COLOR

    # Slide 7: Clean Code & Architecture Optimization
    slide = prs.slides.add_slide(slide_layout)
    add_slide_title(slide, "Refactoring & Workspace Clean-up")
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    refactors = [
        ("Removal of Disconnected AI Service", "The legacy Python FastAPI crop disease service folder (`ai-service`) was completely unlinked from operational code and has been removed to conserve workspace resources."),
        ("Deletion of Unused Image Upload Code", "Eliminated unused `uploadMiddleware.js` (Multer setup) and `cloudinary.js` helpers in the backend since image uploads were not active, reducing server bundle overhead."),
        ("Package Clean-up", "Removed unused `cloudinary` and `multer` dependencies from `server/package.json` to keep development and production builds lean and secure."),
        ("Database Clean-up", "Removed unnecessary database seeding for Mandi Market Prices in `server.js` as requested to align the backend configuration with current operational needs.")
    ]
    
    for i, (title, desc) in enumerate(refactors):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run_title = p.add_run()
        run_title.text = f"•  {title}: "
        run_title.font.bold = True
        run_title.font.size = Pt(18)
        run_title.font.color.rgb = SECONDARY_COLOR
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(16)
        run_desc.font.color.rgb = TEXT_COLOR

    # Slide 8: Future Extensions
    slide = prs.slides.add_slide(slide_layout)
    add_slide_title(slide, "Future Roadmap")
    
    content_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(4.8))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    roadmap = [
        ("Integrated AI Diagnostic Engine", "Re-integrate a production-ready plant disease classification tool using TensorFlow Lite on the backend/edge, complete with a frontend camera upload interface."),
        ("Offline-First Mobile Support", "Optimize the client SPA as a Progressive Web App (PWA) with Service Workers to cache advisories and schemes, assisting farmers even in areas with weak cellular networks."),
        ("Localized Audio/Voice Advisories", "Utilize Speech Synthesis (Text-to-Speech) APIs to read out weather advisories and scheme steps, supporting low-literacy users."),
        ("Mandi Market Analytics & Pricing Forecasts", "Integrate active government APIs (e.g., data.gov.in) to retrieve live, non-simulated local market pricing graphs with future predictive trends.")
    ]
    
    for i, (title, desc) in enumerate(roadmap):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        run_title = p.add_run()
        run_title.text = f"•  {title}: "
        run_title.font.bold = True
        run_title.font.size = Pt(18)
        run_title.font.color.rgb = PRIMARY_COLOR
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(16)
        run_desc.font.color.rgb = TEXT_COLOR

    # Slide 9: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.833), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Thank You!"
    p.font.name = 'Arial'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "KrishiDrishti – Intelligent Farming Assistant"
    p2.font.name = 'Arial'
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = DARK_COLOR
    p2.space_before = Pt(10)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Questions & Discussions"
    p3.font.name = 'Arial'
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_COLOR
    p3.space_before = Pt(20)
    p3.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save('KrishiDrishti_Presentation.pptx')
    print("Presentation created successfully as 'KrishiDrishti_Presentation.pptx'!")

if __name__ == '__main__':
    create_presentation()
